"""Smoke tests for the Phase 1 bait PPTX renderer.

Verifies the renderer runs end-to-end on a small synthetic Phase1Report,
produces a file that python-pptx can round-trip, and that key text content
ends up in the slide. We do not assert exact byte equality — layout is
deterministic in text content but python-pptx zip internals (timestamps,
shape ids) can vary.
"""

from __future__ import annotations

import zipfile
from datetime import date, datetime, timezone
from pathlib import Path

import pytest

# Gate the whole module on python-pptx being available. The main phase1
# pipeline does not depend on it; only the bait renderer does.
pptx = pytest.importorskip("pptx")

from src.voc.reporting.phase1.bait import render_bait_pptx
from src.voc.reporting.phase1.schema import (
    ChannelSignals,
    DeterministicMetrics,
    NarrativeBlock,
    Phase1Report,
    ProductInScope,
    RatingMetrics,
    ReportProvenance,
    ReportQuery,
    ReportScope,
    SegmentMetrics,
    SignalCandidate,
    SignalsBundle,
    TimeWindow,
)


def _tiny_report() -> Phase1Report:
    """A minimally-populated Phase1Report that exercises every bait section.

    Includes: 2 cautionary signals, 1 gap rule (authenticity → high-severity
    promotion), realistic KPI values, one product in scope with a label.
    """
    return Phase1Report(
        report_id="bait_smoke_test",
        generated_at=datetime(2026, 4, 23, 12, 0, tzinfo=timezone.utc),
        query=ReportQuery(),
        scope=ReportScope(
            channels=["oliveyoung"],
            products=[ProductInScope(
                product_id="P1", channel="oliveyoung",
                display_label="Test Product", n_reviews=100,
            )],
            total_reviews=100,
        ),
        deterministic_metrics=DeterministicMetrics(
            total_reviews=100,
            n_products=1,
            channels={"oliveyoung": 100},
            rating=RatingMetrics(
                n=100, missing=0,
                avg_raw=4.5,
                distribution_raw={5: 70, 4: 20, 3: 5, 2: 3, 1: 2},
            ),
            time_window=TimeWindow(
                start_date=date(2026, 1, 1),
                end_date=date(2026, 4, 1),
                days_span=90,
            ),
            dominant_product=None,
            per_product=[],
            segments=SegmentMetrics(),
            channel_signals=ChannelSignals(),
        ),
        signals=SignalsBundle(
            positive=[],
            cautionary=[
                SignalCandidate(
                    name="packaging_complaint",
                    display_label="포장·개봉·배송 파손 불만",
                    category="cautionary",
                    evidence_count=5,
                    coverage_ratio=0.05,
                    sample_review_ids=["r_pkg_1", "r_pkg_2", "r_pkg_3"],
                ),
                SignalCandidate(
                    name="pigment_complaint",
                    display_label="발색 불만",
                    category="cautionary",
                    evidence_count=3,
                    coverage_ratio=0.03,
                    sample_review_ids=["r_pig_1"],
                ),
            ],
            gaps=[
                SignalCandidate(
                    name="coupang_authenticity_concern",
                    display_label="정품·가품 의심 언급 (고위험 운영 신호)",
                    category="gap",
                    evidence_count=1,
                    coverage_ratio=0.01,
                    sample_review_ids=["r_auth_1"],
                ),
            ],
        ),
        narrative=NarrativeBlock(
            summary_md="", sections_md={}, caveats=[], source="template",
        ),
        provenance=ReportProvenance(
            phase1_run_ids=["test_run"],
            sample_review_ids=[],
            lexicon_version="positive=test;cautionary=test",
            llm_model=None,
            llm_prompt_hash=None,
        ),
    )


def _collect_slide_text(pptx_path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs).strip()
                if text:
                    out.append(text)
    return "\n".join(out)


class TestBaitSmoke:
    def test_renders_valid_pptx(self, tmp_path: Path) -> None:
        report = _tiny_report()
        out = render_bait_pptx(report, tmp_path / "smoke.pptx")
        assert out.is_file()
        assert zipfile.is_zipfile(out)
        assert out.stat().st_size > 10 * 1024

    def test_slide_contains_expected_sections(self, tmp_path: Path) -> None:
        report = _tiny_report()
        out = render_bait_pptx(report, tmp_path / "smoke.pptx")
        text = _collect_slide_text(out)
        assert "Test Product" in text
        assert "4.50" in text
        assert "100" in text
        assert "oliveyoung" in text
        assert "관찰된 주요 주의 신호" in text
        assert "운영 관찰" in text
        assert "해석" in text
        assert "추천 점검 항목" in text
        assert "포장·개봉·배송 파손 불만" in text
        assert "발색 불만" in text
        assert "정품·가품 의심" in text
        assert "규칙 기반" in text

    def test_runs_without_review_texts(self, tmp_path: Path) -> None:
        report = _tiny_report()
        out = render_bait_pptx(
            report, tmp_path / "no_quotes.pptx",
            review_text_by_id=None,
        )
        assert zipfile.is_zipfile(out)
        text = _collect_slide_text(out)
        assert "포장·개봉·배송 파손 불만" in text
        assert "“" not in text

    def test_renders_quoted_excerpts_when_lookup_provided(
        self, tmp_path: Path,
    ) -> None:
        report = _tiny_report()
        quote_map = {
            "r_pkg_1": "뚜껑이 자꾸 열려요 포장이 부실합니다",
            "r_pig_1": "발색이 기대 이하입니다",
            "r_auth_1": "정품이 아닌 것 같아 걱정됩니다",
        }
        out = render_bait_pptx(
            report, tmp_path / "with_quotes.pptx",
            review_text_by_id=quote_map,
        )
        text = _collect_slide_text(out)
        assert "뚜껑이 자꾸 열려요" in text
        assert "발색이 기대 이하" in text
        assert "정품이 아닌 것 같" in text

    def test_empty_signals_renders_graceful_state(self, tmp_path: Path) -> None:
        report = _tiny_report()
        report.signals = SignalsBundle(positive=[], cautionary=[], gaps=[])
        out = render_bait_pptx(report, tmp_path / "empty.pptx")
        assert zipfile.is_zipfile(out)
        text = _collect_slide_text(out)
        assert "Test Product" in text
        assert "반복 주의 신호가 도출되지 않았습니다" in text
