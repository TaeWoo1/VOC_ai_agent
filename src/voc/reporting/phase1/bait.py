"""1-page outbound 'bait report' renderer for Phase 1 VOC artifacts.

Strict post-stage renderer: consumes an already-built ``Phase1Report`` (the
evidence-layer JSON from ``scripts/generate_phase1_report.py``) and produces
a single-slide PPTX that a seller-facing recipient could read in 15 seconds.

The existing JSON + Markdown report remains the evidence layer. This layer
is deliberately lossy and scan-optimised: KPI cards, top 3 issues with
quoted evidence, one operational observation, a cautious template
interpretation, 2–3 recommended checks, and a disclaimer footer. No charts,
no causal claims, no LLM wording.

Install the dependency: ``pip install -e '.[bait]'``. The rest of the
Phase 1 stack does NOT depend on python-pptx.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Iterable

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu, Inches, Pt
except ImportError as e:  # pragma: no cover
    raise ImportError(
        "bait.py requires python-pptx. Install with: pip install -e '.[bait]'"
    ) from e

from src.voc.reporting.phase1.schema import (
    Phase1Report,
    SignalCandidate,
    SignalsBundle,
)


# ---------------------------------------------------------------------------
# Layout constants — single A4-ratio landscape slide, all coordinates in EMU
# via python-pptx helpers. Tuned by eye; adjust in place if the rendered
# artifact overflows a given block on real data.
# ---------------------------------------------------------------------------

# Slide: use widescreen 13.333" x 7.5" (standard 16:9 for modern decks)
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)

# Colour palette — restrained, single accent
_COLOR_TITLE = RGBColor(0x1A, 0x1A, 0x1A)       # near-black headings
_COLOR_BODY = RGBColor(0x33, 0x33, 0x33)         # body text
_COLOR_MUTED = RGBColor(0x77, 0x77, 0x77)        # captions / disclaimer
_COLOR_ACCENT = RGBColor(0x1F, 0x4E, 0x79)       # navy-blue accent for KPI / numbers
_COLOR_CARD_BG = RGBColor(0xF4, 0xF6, 0xF8)      # light-grey card fill
_COLOR_CARD_BORDER = RGBColor(0xDC, 0xE0, 0xE5)
_COLOR_QUOTE = RGBColor(0x55, 0x55, 0x55)
_COLOR_OPERATIONAL = RGBColor(0xFB, 0xF3, 0xEA)  # pale-amber for operational observation

# High-severity gaps surfaced in the operational block (same policy as
# narrative.py's exec summary). Coupang-authenticity and skin-irritation
# always promote; repurchase-mismatch promotes only when coverage is
# material.
_HIGH_SEVERITY_GAP_NAMES = frozenset({
    "coupang_authenticity_concern",
    "skin_irritation_concern",
})
_REPURCHASE_MISMATCH_NAME = "api_repurchase_vs_text_mention"
_REPURCHASE_SIGNIFICANT_PCT = 0.05

# Quote-excerpt character cap for slide legibility (stricter than the
# narrative's 160 — PPTX cells don't wrap as cleanly).
_QUOTE_MAX_CHARS = 110


# ---------------------------------------------------------------------------
# Suggested-check phrasings — conservative, action-framed, no causal claims.
# Keyed by signal_id; missing ids produce no suggestion. Order of insertion
# into the slide follows the signal's evidence-count ranking.
# ---------------------------------------------------------------------------

_SUGGESTION_BY_SIGNAL: dict[str, str] = {
    "packaging_complaint": "포장 품질과 배송 파손 이슈를 운영팀과 함께 점검 권장",
    "pigment_complaint": "대표 이미지와 실제 셰이드의 일치 여부 확인 권장",
    "shade_mismatch": "상품 상세 이미지·셰이드 설명의 정확도 재검토 권장",
    "application_issue": "권장 사용법·바르기 가이드 커뮤니케이션 보완 검토",
    "tone_mismatch": "퍼스널컬러 관련 상품 설명 위치·표현 점검 권장",
    "value_complaint": "가격·가치 포지셔닝 재검토 권장",
    "persistence_reservation": "지속력 관련 사용자 기대치 커뮤니케이션 점검",
    "api_repurchase_vs_text_mention": "재구매 API 귀속 로직 정확성 감사 권장",
    "coupang_authenticity_concern": "쿠팡 셀러 리스팅의 정품·가품 여부 점검 필요",
    "skin_irritation_concern": "피부 자극 관련 사례를 안전 담당 부서와 공유 권장",
    "eyeshadow_fallout": "팔레트 가루 날림 이슈를 제조·QA와 함께 점검 권장",
}


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------


def render_bait_pptx(
    report: Phase1Report,
    output_path: Path | str,
    *,
    review_text_by_id: dict[str, str] | None = None,
) -> Path:
    """Render a single-slide bait PPTX from a ``Phase1Report``.

    ``review_text_by_id``: optional ``{review_id: text}`` lookup. When
    provided, quoted excerpts render in the top-issue and operational
    blocks; missing ids silently skip. When omitted, those blocks render
    with the signal name and counts only — still usable, just less
    evidence-rich. The caller (CLI) is responsible for producing this map
    via a DB lookup on the report's sample_review_ids.

    Returns the written path.
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    blank_layout = prs.slide_layouts[6]  # fully blank
    slide = prs.slides.add_slide(blank_layout)

    # Cursor y-position, updated as each section lays out top-to-bottom.
    cur_y = Inches(0.25)

    cur_y = _render_title_strip(slide, report, left=Inches(0.4),
                                top=cur_y, width=Inches(12.5))
    cur_y += Inches(0.15)

    cur_y = _render_kpi_row(slide, report, left=Inches(0.4),
                            top=cur_y, width=Inches(12.5), height=Inches(0.9))
    cur_y += Inches(0.25)

    cur_y = _render_top_issues(slide, report, review_text_by_id,
                               left=Inches(0.4), top=cur_y,
                               width=Inches(12.5), height=Inches(1.7))
    cur_y += Inches(0.15)

    cur_y = _render_operational(slide, report, review_text_by_id,
                                left=Inches(0.4), top=cur_y,
                                width=Inches(12.5))
    cur_y += Inches(0.15)

    cur_y = _render_interpretation(slide, report, left=Inches(0.4),
                                   top=cur_y, width=Inches(12.5))
    cur_y += Inches(0.1)

    cur_y = _render_suggestions(slide, report, left=Inches(0.4),
                                top=cur_y, width=Inches(12.5))

    # Footer disclaimer pinned to the bottom of the slide regardless of
    # where the content cursor landed.
    _render_footer(slide, report, left=Inches(0.4),
                   top=Inches(7.05), width=Inches(12.5))

    prs.save(str(output_path))
    return output_path


# ---------------------------------------------------------------------------
# Section renderers
# ---------------------------------------------------------------------------


def _render_title_strip(slide, report: Phase1Report, *,
                        left, top, width) -> Any:
    """Product name / window / generated-date row."""
    product_name = _title_subject(report)
    window = report.deterministic_metrics.time_window
    window_str = ""
    if window.start_date and window.end_date:
        window_str = f"{window.start_date.isoformat()} ~ {window.end_date.isoformat()}"

    gen_date = report.generated_at.date().isoformat() if report.generated_at else ""

    box = slide.shapes.add_textbox(left, top, width, Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    # Line 1: product name + window (larger, bolder)
    p = tf.paragraphs[0]
    run_title = p.add_run()
    run_title.text = f"{product_name}"
    run_title.font.size = Pt(20)
    run_title.font.bold = True
    run_title.font.color.rgb = _COLOR_TITLE

    if window_str:
        run_window = p.add_run()
        run_window.text = f"   ·   {window_str}"
        run_window.font.size = Pt(12)
        run_window.font.color.rgb = _COLOR_MUTED
        run_window.font.bold = False

    # Line 2: generated date (small, muted)
    p2 = tf.add_paragraph()
    run_gen = p2.add_run()
    run_gen.text = f"Generated {gen_date}"
    run_gen.font.size = Pt(9)
    run_gen.font.color.rgb = _COLOR_MUTED

    return top + Inches(0.55)


def _render_kpi_row(slide, report: Phase1Report, *,
                    left, top, width, height) -> Any:
    """Four KPI cards: avg rating, 5★ share, review count, channels."""
    m = report.deterministic_metrics
    dist = m.rating.distribution_raw
    rated_total = sum(dist.values())
    five_star_count = dist.get(5, 0)
    five_pct = (100 * five_star_count / rated_total) if rated_total else 0.0

    avg_rating_str = (
        f"{round(m.rating.avg_raw, 2):.2f}" if m.rating.avg_raw is not None else "—"
    )

    cards = [
        ("평균 평점", avg_rating_str, "/ 5"),
        ("5★ 비중", f"{round(five_pct, 1):g}%", ""),
        ("분석 리뷰", f"{m.total_reviews}", "건"),
        ("채널", "+".join(sorted(m.channels)) or "—", ""),
    ]

    n_cards = len(cards)
    gap = Inches(0.15)
    card_w = Emu((int(width) - int(gap) * (n_cards - 1)) // n_cards)

    for i, (label, value, unit) in enumerate(cards):
        x = Emu(int(left) + i * (int(card_w) + int(gap)))
        _draw_card(slide, x, top, card_w, height, fill=_COLOR_CARD_BG,
                   border=_COLOR_CARD_BORDER)

        # Label (small, muted)
        label_box = slide.shapes.add_textbox(
            x + Inches(0.15), top + Inches(0.1),
            card_w - Inches(0.3), Inches(0.25),
        )
        lf = label_box.text_frame
        lf.margin_top = Emu(0)
        lf.margin_bottom = Emu(0)
        p = lf.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.size = Pt(9)
        r.font.color.rgb = _COLOR_MUTED

        # Value (big, accent) + unit (small, muted)
        val_box = slide.shapes.add_textbox(
            x + Inches(0.15), top + Inches(0.35),
            card_w - Inches(0.3), Inches(0.55),
        )
        vf = val_box.text_frame
        vf.margin_top = Emu(0)
        vf.margin_bottom = Emu(0)
        pv = vf.paragraphs[0]
        rv = pv.add_run()
        rv.text = value
        rv.font.size = Pt(22)
        rv.font.bold = True
        rv.font.color.rgb = _COLOR_ACCENT
        if unit:
            ru = pv.add_run()
            ru.text = f" {unit}"
            ru.font.size = Pt(11)
            ru.font.color.rgb = _COLOR_MUTED

    return top + height


def _render_top_issues(slide, report: Phase1Report,
                       review_text_by_id: dict[str, str] | None, *,
                       left, top, width, height) -> Any:
    """Top 3 cautionary signals in 3 side-by-side cards."""
    # Section header
    hdr_box = slide.shapes.add_textbox(left, top, width, Inches(0.25))
    hf = hdr_box.text_frame
    hf.margin_top = Emu(0)
    hf.margin_bottom = Emu(0)
    hp = hf.paragraphs[0]
    hr = hp.add_run()
    hr.text = "관찰된 주요 주의 신호 (상위 3)"
    hr.font.size = Pt(11)
    hr.font.bold = True
    hr.font.color.rgb = _COLOR_TITLE

    cur_top = top + Inches(0.3)

    cautionary_sorted = sorted(
        report.signals.cautionary,
        key=lambda s: s.evidence_count, reverse=True,
    )
    top3 = cautionary_sorted[:3]

    card_h = height - Inches(0.3)
    if not top3:
        # Graceful empty state: one wide block with a cautious message.
        _render_empty_block(slide, left, cur_top, width, card_h,
                            text="규칙 기반 분석에서 반복 주의 신호가 도출되지 않았습니다.")
        return cur_top + card_h

    gap = Inches(0.2)
    n = len(top3)
    card_w = Emu((int(width) - int(gap) * (n - 1)) // n)

    for i, sig in enumerate(top3):
        x = Emu(int(left) + i * (int(card_w) + int(gap)))
        _draw_card(slide, x, cur_top, card_w, card_h,
                   fill=_COLOR_CARD_BG, border=_COLOR_CARD_BORDER)

        tb = slide.shapes.add_textbox(
            x + Inches(0.15), cur_top + Inches(0.1),
            card_w - Inches(0.3), card_h - Inches(0.2),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)

        # Title row: rank + label
        p = tf.paragraphs[0]
        rank = p.add_run()
        rank.text = f"#{i + 1}   "
        rank.font.size = Pt(10)
        rank.font.bold = True
        rank.font.color.rgb = _COLOR_ACCENT

        lbl = p.add_run()
        lbl.text = sig.display_label
        lbl.font.size = Pt(11)
        lbl.font.bold = True
        lbl.font.color.rgb = _COLOR_TITLE

        # Counts row
        p2 = tf.add_paragraph()
        r = p2.add_run()
        r.text = f"{sig.evidence_count}건 · {_fmt_pct(sig.coverage_ratio)}"
        r.font.size = Pt(9)
        r.font.color.rgb = _COLOR_MUTED

        # Quote (if available)
        quote = _pick_quote(sig, review_text_by_id)
        if quote:
            p3 = tf.add_paragraph()
            rq = p3.add_run()
            rq.text = f"  “{quote}”"
            rq.font.size = Pt(9)
            rq.font.italic = True
            rq.font.color.rgb = _COLOR_QUOTE

    return cur_top + card_h


def _render_operational(slide, report: Phase1Report,
                        review_text_by_id: dict[str, str] | None, *,
                        left, top, width) -> Any:
    """One operational-observation card surfacing the most salient gap rule.

    Selection priority: high-severity gap rules first (authenticity, skin
    irritation), then repurchase-mismatch if coverage ≥ 5%, else highest-
    coverage gap rule, else skip.
    """
    op = _pick_operational_signal(report.signals)
    if op is None:
        return top  # nothing to render; don't reserve space

    card_h = Inches(0.75)

    # Header line (small, pre-card)
    hdr_box = slide.shapes.add_textbox(left, top, width, Inches(0.22))
    hf = hdr_box.text_frame
    hf.margin_top = Emu(0)
    hf.margin_bottom = Emu(0)
    hp = hf.paragraphs[0]
    hr = hp.add_run()
    hr.text = "운영 관찰 (주목할 신호)"
    hr.font.size = Pt(11)
    hr.font.bold = True
    hr.font.color.rgb = _COLOR_TITLE

    card_top = top + Inches(0.25)
    _draw_card(slide, left, card_top, width, card_h,
               fill=_COLOR_OPERATIONAL, border=_COLOR_CARD_BORDER)

    tb = slide.shapes.add_textbox(
        left + Inches(0.2), card_top + Inches(0.08),
        width - Inches(0.4), card_h - Inches(0.15),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"{op.display_label} — {op.evidence_count}건 ({_fmt_pct(op.coverage_ratio)})"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = _COLOR_TITLE

    quote = _pick_quote(op, review_text_by_id)
    if quote:
        p2 = tf.add_paragraph()
        rq = p2.add_run()
        rq.text = f"  “{quote}”"
        rq.font.size = Pt(9)
        rq.font.italic = True
        rq.font.color.rgb = _COLOR_QUOTE

    return card_top + card_h


def _render_interpretation(slide, report: Phase1Report, *,
                           left, top, width) -> Any:
    """Cautious 1–3 sentence interpretation. Pattern observation only."""
    text = _compose_interpretation(report)
    if not text:
        return top

    box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    # Mini-header
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "해석 (관찰된 패턴)"
    r0.font.size = Pt(11)
    r0.font.bold = True
    r0.font.color.rgb = _COLOR_TITLE

    p = tf.add_paragraph()
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.color.rgb = _COLOR_BODY

    return top + Inches(0.6)


def _render_suggestions(slide, report: Phase1Report, *,
                        left, top, width) -> Any:
    """Up to 3 recommended-check bullets drawn from fired cautionary +
    operational signals. Conservative verbs throughout."""
    bullets = _compose_suggestions(report)
    if not bullets:
        return top

    box = slide.shapes.add_textbox(left, top, width, Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "추천 점검 항목 (해석 기반)"
    r0.font.size = Pt(11)
    r0.font.bold = True
    r0.font.color.rgb = _COLOR_TITLE

    for b in bullets:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"•  {b}"
        r.font.size = Pt(10)
        r.font.color.rgb = _COLOR_BODY

    return top + Inches(0.85)


def _render_footer(slide, report: Phase1Report, *,
                   left, top, width) -> None:
    """Disclaimer footer: scope, method framing, versions."""
    m = report.deterministic_metrics
    window = m.time_window
    window_part = ""
    if window.start_date and window.end_date:
        window_part = (
            f" · {window.start_date.isoformat()}~{window.end_date.isoformat()}"
        )
    channels_part = "+".join(sorted(m.channels)) if m.channels else "—"
    lex_part = f" · lexicon {report.provenance.lexicon_version}"

    text = (
        f"분석 기준: 리뷰 {m.total_reviews}건 · 채널 {channels_part}"
        f"{window_part}{lex_part}. "
        "규칙 기반 통계적 패턴 관찰이며 인과적 진단이 아닙니다. "
        "인용문은 대표 사례이며 전수가 아닙니다."
    )

    box = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(8)
    r.font.color.rgb = _COLOR_MUTED


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _draw_card(slide, left, top, width, height, *, fill, border) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.5)
    # Don't attach text; text boxes overlay on top for precise control.
    shape.text_frame.text = ""
    shape.shadow.inherit = False


def _render_empty_block(slide, left, top, width, height, *, text: str) -> None:
    _draw_card(slide, left, top, width, height,
               fill=_COLOR_CARD_BG, border=_COLOR_CARD_BORDER)
    tb = slide.shapes.add_textbox(left + Inches(0.2),
                                  top + Inches(0.1),
                                  width - Inches(0.4),
                                  height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = _COLOR_MUTED


def _title_subject(report: Phase1Report) -> str:
    dp = report.deterministic_metrics.dominant_product
    if dp is not None:
        label = next(
            (p.display_label for p in report.scope.products
             if p.product_id == dp.product_id and p.display_label),
            None,
        )
        return label or dp.product_id
    # No dominant_product computed — fall back to the first scoped product
    # so single-product reports still surface a real title. Only the truly
    # empty case drops to the generic label.
    if report.scope.products:
        first = report.scope.products[0]
        return first.display_label or first.product_id
    return "Phase 1 VOC 리포트"


def _fmt_pct(x: float) -> str:
    return f"{round(x * 100, 1):g}%"


def _pick_quote(sig: SignalCandidate,
                review_text_by_id: dict[str, str] | None) -> str | None:
    if not review_text_by_id or not sig.sample_review_ids:
        return None
    for rid in sig.sample_review_ids:
        text = review_text_by_id.get(rid)
        if not text:
            continue
        snippet = " ".join(text.split())
        if not snippet:
            continue
        if len(snippet) > _QUOTE_MAX_CHARS:
            snippet = snippet[:_QUOTE_MAX_CHARS].rstrip() + "…"
        return snippet
    return None


def _pick_operational_signal(signals: SignalsBundle) -> SignalCandidate | None:
    """Select the single most-outbound-interesting gap rule, or None."""
    gaps = signals.gaps
    if not gaps:
        return None
    # 1. High-severity first
    for g in gaps:
        if g.name in _HIGH_SEVERITY_GAP_NAMES:
            return g
    # 2. Repurchase mismatch if material
    for g in gaps:
        if (g.name == _REPURCHASE_MISMATCH_NAME
                and g.coverage_ratio >= _REPURCHASE_SIGNIFICANT_PCT):
            return g
    # 3. Highest-coverage gap otherwise (cautiously — might be noise)
    return max(gaps, key=lambda g: g.coverage_ratio, default=None)


def _compose_interpretation(report: Phase1Report) -> str:
    """Produce a 1–3 sentence cautious interpretation. Template only.

    Rules:
      - Rating framed as 'observed posture', never 'excellent' / 'weak'.
      - Cautionary concentration framed as 'concentrated on' / 'observed on',
        never causal.
      - Operational caveat included only when an operational signal was
        surfaced, and framed as 'suggests attention', never 'indicates a
        problem'.
    """
    m = report.deterministic_metrics
    s = report.signals
    parts: list[str] = []

    if m.total_reviews == 0:
        return "분석 대상 리뷰가 없어 관찰 결과 없음."

    if m.rating.avg_raw is not None:
        dist = m.rating.distribution_raw
        rated = sum(dist.values())
        five_pct = (100 * dist.get(5, 0) / rated) if rated else 0.0
        parts.append(
            f"평균 평점 {round(m.rating.avg_raw, 2):.2f}/5, "
            f"5★ 비중 {round(five_pct, 1):g}%로 전반 긍정도는 높은 편으로 관찰됩니다."
        )

    top_c = sorted(s.cautionary, key=lambda x: x.evidence_count, reverse=True)
    if len(top_c) >= 2:
        parts.append(
            f"주의 신호는 '{top_c[0].display_label}'와(과) "
            f"'{top_c[1].display_label}'에 상대적으로 집중된 패턴으로 관찰됩니다."
        )
    elif len(top_c) == 1:
        parts.append(
            f"주의 신호는 '{top_c[0].display_label}'에 집중된 패턴으로 관찰됩니다."
        )

    op = _pick_operational_signal(s)
    if op is not None:
        parts.append(
            f"'{op.display_label}'은 해석·추가 확인이 필요한 지점으로 보입니다."
        )

    return " ".join(parts)


def _compose_suggestions(report: Phase1Report) -> list[str]:
    """Up to 3 hedged action items keyed off fired signals."""
    seen: set[str] = set()
    out: list[str] = []

    # Rank cautionary by evidence count; append operational pick if any.
    cautionary_sorted = sorted(
        report.signals.cautionary,
        key=lambda s: s.evidence_count, reverse=True,
    )
    op = _pick_operational_signal(report.signals)

    candidates: Iterable[SignalCandidate] = list(cautionary_sorted)
    if op is not None:
        candidates = [op] + list(cautionary_sorted)

    for sig in candidates:
        text = _SUGGESTION_BY_SIGNAL.get(sig.name)
        if text is None:
            continue
        if text in seen:
            continue
        seen.add(text)
        out.append(text)
        if len(out) >= 3:
            break
    return out
